#!/usr/bin/env python3
"""
generate_deck_v3.py — Rebuild jboss-ai-monitor-deck.pptx (12 slides)
Follows the official Red Hat 2025 presentation template design system.

Canvas: 13.33" x 7.50" (standard 16:9)
Colors: #EE0000 red, #5E40BE purple-50, #21134D purple-70 (dark bg),
        black, white, grays
Layout: section label (0.49,0.06), red rule @y=0.95, title (0.97,1.2)
        at 28pt, body 18pt, slide number bottom-left (0.09,6.75)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Color palette (RH 2025) ───────────────────────────────────────────────────
RED        = RGBColor(0xEE, 0x00, 0x00)   # red-50
PURPLE     = RGBColor(0x5E, 0x40, 0xBE)   # purple-50 — section divider text
DARK_PUR   = RGBColor(0x21, 0x13, 0x4D)   # purple-70 — dark section bg
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY80     = RGBColor(0x29, 0x29, 0x29)
GRAY60     = RGBColor(0x4D, 0x4D, 0x4D)
GRAY40     = RGBColor(0xA3, 0xA3, 0xA3)
GRAY20     = RGBColor(0xE0, 0xE0, 0xE0)
TEAL       = RGBColor(0x37, 0xA3, 0xA3)   # teal-50 — accent callouts
GREEN      = RGBColor(0x3E, 0x8C, 0x35)   # positive indicators

# ── Slide dimensions ──────────────────────────────────────────────────────────
W = 13.33
H = 7.50


# ── Low-level helpers ─────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size,
             bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Red Hat Display" if bold else "Red Hat Text"
    return tb


def add_para(tf, text, font_size, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    from pptx.util import Pt as Pt2
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt2(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt2(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Red Hat Display" if bold else "Red Hat Text"
    return p


def bullet_para(tf, text, font_size=14, color=BLACK, indent_level=0):
    p = tf.add_paragraph()
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = "Red Hat Text"
    return p


# ── Slide chrome (shared across all content slides) ──────────────────────────

def add_content_chrome(slide, section_label, slide_num):
    """Add section label, red rule, and slide number."""
    # Section label top-left
    add_text(slide, 0.49, 0.06, 5.63, 0.35,
             section_label, 9, color=GRAY60)
    # Red horizontal rule below label
    add_rect(slide, 0.49, 0.52, W - 0.98, 0.04, RED)
    # Slide number bottom-left
    add_text(slide, 0.09, 6.75, 0.8, 0.25,
             str(slide_num), 9, color=GRAY60)


def add_section_divider(slide, section_label, title, slide_num):
    """Dark purple section-divider slide (RH 2025 template section pattern)."""
    # Dark background
    add_rect(slide, 0, 0, W, H, DARK_PUR)
    # Small section label
    add_text(slide, 1.5, 1.8, 9.0, 0.6,
             section_label, 18, color=PURPLE)
    # Large title
    tb = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.4), Inches(10.0), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(54)
    run.font.color.rgb = PURPLE
    run.font.name = "Red Hat Display"
    # Slide number
    add_text(slide, 0.09, 6.75, 0.8, 0.25,
             str(slide_num), 9, color=GRAY40)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(slide):
    """Title slide — dark background, red hat branding."""
    # Full dark bg
    add_rect(slide, 0, 0, W, H, DARK_PUR)
    # Red accent bar left edge
    add_rect(slide, 0, 0, 0.18, H, RED)
    # Eyebrow
    add_text(slide, 0.5, 1.6, 11.0, 0.5,
             "Red Hat OpenShift AI  ·  JBoss/WildFly  ·  OpenShift", 13,
             color=PURPLE)
    # Main title
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2), Inches(11.5), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Autonomous JBoss Operations"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Red Hat Display"
    add_para(tf, "with Red Hat OpenShift AI", 30, color=RED)
    # Subtitle
    add_text(slide, 0.5, 4.9, 9.0, 0.7,
             "Detect. Analyse. Resolve. — Automatically.", 18,
             italic=True, color=GRAY40)
    # Presenter line
    add_text(slide, 0.5, 5.7, 9.0, 0.5,
             "Nilay Saraiya  ·  Senior Architect, Red Hat", 12, color=GRAY40)
    # Red Hat wordmark area bottom right
    add_rect(slide, 11.2, 6.4, 1.93, 0.6, RED)
    add_text(slide, 11.25, 6.45, 1.85, 0.5,
             "RED HAT", 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_02_problem(slide):
    add_content_chrome(slide, "The Challenge", 2)
    # Title
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "The Status Quo: Manual, Reactive, Expensive", 22,
             bold=True, color=BLACK)
    # Subheading
    add_text(slide, 0.97, 1.4, 11.4, 0.45,
             "Every JBoss incident today triggers a costly, error-prone manual response", 14,
             italic=True, color=GRAY60)
    # Red rule under title
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    # Three pain-point cards
    cards = [
        ("2–4 hours", "per incident", "Senior engineer time spent triage, diagnosis, and ticket writing"),
        ("3am alerts", "every incident", "On-call engineers paged for issues that follow predictable patterns"),
        ("Inconsistent\ntickets", "every time", "Different engineers write different tickets — quality varies widely"),
    ]
    cols = [0.6, 4.9, 9.2]
    for (stat, label, desc), x in zip(cards, cols):
        add_rect(slide, x, 2.1, 3.5, 4.4, GRAY20)
        add_rect(slide, x, 2.1, 3.5, 0.08, RED)
        add_text(slide, x+0.15, 2.3, 3.2, 0.9,
                 stat, 30, bold=True, color=RED, align=PP_ALIGN.CENTER)
        add_text(slide, x+0.15, 3.25, 3.2, 0.4,
                 label, 12, color=GRAY60, align=PP_ALIGN.CENTER)
        add_rect(slide, x+0.15, 3.7, 3.2, 0.04, GRAY40)
        add_text(slide, x+0.15, 3.85, 3.2, 2.5,
                 desc, 13, color=BLACK, wrap=True)


def slide_03_solution(slide):
    add_content_chrome(slide, "The Solution", 3)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "Detect. Analyse. Document. Automatically.", 22,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 9.0, 0.4,
             "A closed-loop AI agent running inside OpenShift", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    # Architecture flow — 4 boxes with arrows
    steps = [
        ("Monitor", "4 detection layers:\nPod crashes · Log errors\nAlerts · Health checks"),
        ("Deduplicate", "MD5 fingerprint\nsuppresses repeats\nfor 2 hours"),
        ("Analyse\n(RHOAI LLM)", "llama-3.2-3B via vLLM\nRoot cause + steps\nvia function-calling"),
        ("JIRA Ticket", "Structured ADF ticket\nwith root cause,\nsteps & prevention"),
    ]
    colors = [DARK_PUR, GRAY60, RED, TEAL]
    for i, ((title, body), x, col) in enumerate(
            zip(steps, [0.6, 3.55, 6.5, 9.45], colors)):
        add_rect(slide, x, 2.1, 2.9, 4.2, col)
        add_text(slide, x+0.1, 2.25, 2.7, 0.7,
                 title, 16, bold=True, color=WHITE, wrap=True)
        add_text(slide, x+0.1, 3.0, 2.7, 3.1,
                 body, 13, color=WHITE, wrap=True)
        if i < 3:
            add_text(slide, x+2.88, 3.85, 0.45, 0.5, "→", 22, color=GRAY40)

    # Loop annotation
    add_text(slide, 0.6, 6.45, 12.0, 0.4,
             "Repeats every 60 seconds  ·  Zero human intervention required", 11,
             italic=True, color=GRAY60, align=PP_ALIGN.CENTER)


def slide_04_roi(slide):
    add_content_chrome(slide, "Business Value", 4)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "2.5 Hours of Senior Engineer Time → 15 Minutes of Review", 20,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 11.4, 0.4,
             "Conservative assumptions · 10 incidents/month · 1 JBoss instance", 13,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    # Before / After header
    add_rect(slide, 0.6, 2.1, 5.9, 0.5, GRAY20)
    add_text(slide, 0.6, 2.1, 5.9, 0.5,
             "WITHOUT AI Monitor", 13, bold=True, color=GRAY60, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.8, 2.1, 5.9, 0.5, RED)
    add_text(slide, 6.8, 2.1, 5.9, 0.5,
             "WITH AI Monitor", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    before = [
        ("150 hrs/yr", "Triage & analysis time"),
        ("$30,000/yr", "At $200/hr blended rate"),
        ("3am pages", "Unpredictable on-call load"),
        ("Variable", "Ticket quality & completeness"),
    ]
    after = [
        ("15 min/yr", "Human review only"),
        ("~$1,500/yr", "Estimated operational cost"),
        ("Zero", "Engineer interruptions"),
        ("Consistent", "Structured AI-generated tickets"),
    ]
    y = 2.75
    for (bs, bl), (as_, al) in zip(before, after):
        add_rect(slide, 0.6, y, 5.9, 0.72, GRAY20)
        add_text(slide, 0.75, y+0.04, 2.0, 0.35, bs, 18, bold=True, color=GRAY60)
        add_text(slide, 2.85, y+0.1, 3.5, 0.55, bl, 11, color=GRAY60)
        add_rect(slide, 6.8, y, 5.9, 0.72, RGBColor(0xF5, 0xF0, 0xFF))
        add_text(slide, 6.95, y+0.04, 2.0, 0.35, as_, 18, bold=True, color=DARK_PUR)
        add_text(slide, 9.05, y+0.1, 3.5, 0.55, al, 11, color=DARK_PUR)
        add_rect(slide, 0.6, y+0.72, 12.1, 0.04, WHITE)
        y += 0.76

    add_text(slide, 0.6, 6.5, 12.0, 0.35,
             "20× reduction in engineer time  ·  ~$28,500 annual savings per instance",
             12, bold=True, color=RED, align=PP_ALIGN.CENTER)


def slide_05_scale(slide):
    add_content_chrome(slide, "Business Value", 5)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "The Value Scales With Your Environment", 22,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 11.4, 0.4,
             "More instances = more savings, same AI infrastructure cost", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    # Scale table
    headers = ["JBoss Instances", "Incidents/Year", "Engineer Hours Saved", "Annual Value"]
    rows = [
        ["1",   "120",  "148 hrs",   "~$30K"],
        ["5",   "600",  "743 hrs",   "~$149K"],
        ["10",  "1,200","1,485 hrs", "~$297K"],
        ["25",  "3,000","3,713 hrs", "~$743K"],
        ["50+", "6,000+","7,425+ hrs","~$1.5M+"],
    ]
    col_w = [2.5, 2.5, 3.2, 2.8]
    col_x = [0.6, 3.2, 5.8, 9.1]
    # Header row
    hx = 0.6
    for h, cw, cx in zip(headers, col_w, col_x):
        add_rect(slide, cx, 2.1, cw - 0.05, 0.45, DARK_PUR)
        add_text(slide, cx+0.1, 2.15, cw-0.2, 0.35,
                 h, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Data rows
    for ri, row in enumerate(rows):
        bg = GRAY20 if ri % 2 == 0 else WHITE
        highlight = ri == 4
        for val, cw, cx in zip(row, col_w, col_x):
            fill = RGBColor(0xF5, 0xF0, 0xFF) if highlight else bg
            add_rect(slide, cx, 2.6 + ri*0.68, cw-0.05, 0.63, fill)
            tc = RED if (highlight and val.startswith("~$")) else (DARK_PUR if highlight else BLACK)
            add_text(slide, cx+0.1, 2.65+ri*0.68, cw-0.2, 0.55,
                     val, 13, bold=highlight, color=tc, align=PP_ALIGN.CENTER)

    add_text(slide, 0.6, 6.5, 12.0, 0.35,
             "Assumes $200/hr blended rate · 15 min/incident with AI vs 150 min without",
             10, italic=True, color=GRAY60, align=PP_ALIGN.CENTER)


def slide_06_oncall(slide):
    add_content_chrome(slide, "Operational Impact", 6)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "24/7 Coverage — Without 3am Pages", 22,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 9.0, 0.4,
             "AI watches continuously; engineers review in business hours", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    # Left — before timeline
    add_rect(slide, 0.6, 2.1, 5.9, 0.4, GRAY20)
    add_text(slide, 0.6, 2.1, 5.9, 0.4,
             "TODAY — Manual response", 12, bold=True, color=GRAY60, align=PP_ALIGN.CENTER)
    events = [
        ("3:14 AM", "Alert fires — engineer paged"),
        ("3:30 AM", "Engineer logs in, starts triage"),
        ("4:45 AM", "Root cause identified"),
        ("5:15 AM", "JIRA ticket manually written"),
        ("5:30 AM", "Engineer back to sleep (2.25 hrs lost)"),
    ]
    for i, (ts, ev) in enumerate(events):
        y = 2.6 + i * 0.73
        add_rect(slide, 0.7, y, 1.2, 0.55, RED)
        add_text(slide, 0.72, y+0.08, 1.16, 0.4, ts, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 2.0, y+0.1, 4.3, 0.5, ev, 12, color=BLACK)

    # Right — with AI
    add_rect(slide, 6.8, 2.1, 6.3, 0.4, RED)
    add_text(slide, 6.8, 2.1, 6.3, 0.4,
             "WITH AI MONITOR — Zero interruption", 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ai_events = [
        ("3:14 AM", "Crash detected by PodMonitor"),
        ("3:14 AM", "RHOAI analysis complete (< 60s)"),
        ("3:14 AM", "Structured JIRA ticket created"),
        ("9:05 AM", "Engineer reviews ticket — 15 min"),
        ("9:20 AM", "Issue resolved, prevention applied"),
    ]
    for i, (ts, ev) in enumerate(ai_events):
        y = 2.6 + i * 0.73
        add_rect(slide, 6.9, y, 1.2, 0.55, DARK_PUR)
        add_text(slide, 6.92, y+0.08, 1.16, 0.4, ts, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 8.2, y+0.1, 4.7, 0.5, ev, 12, color=BLACK)


def slide_07_ticket(slide):
    add_content_chrome(slide, "Ticket Quality", 7)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "Every Incident Gets a Structured, Actionable Ticket", 20,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 9.0, 0.4,
             "AI-generated tickets follow the same template every time", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    # Simulated JIRA ticket
    add_rect(slide, 0.6, 2.1, 12.1, 4.8, GRAY20)
    add_rect(slide, 0.6, 2.1, 12.1, 0.55, DARK_PUR)
    add_text(slide, 0.75, 2.15, 8.0, 0.45,
             "KAN-42  [HIGH/crash]  OOMKilled: jboss-instance-0", 14,
             bold=True, color=WHITE)
    add_rect(slide, 8.9, 2.2, 1.8, 0.35, RED)
    add_text(slide, 8.9, 2.2, 1.8, 0.35,
             "AI-generated", 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sections = [
        ("Root Cause", "JVM heap exhausted after gradual leak in session cache. "
            "-Xmx1024m limit reached under sustained load. "
            "GC overhead exceeded 98% — OOMKill triggered."),
        ("Resolution Steps",
            "1. Restart pod: oc rollout restart deploy/jboss-instance\n"
            "2. Increase -Xmx to 1536m in JAVA_OPTS_APPEND\n"
            "3. Enable GC logging: -Xlog:gc*:file=/tmp/gc.log\n"
            "4. Review heap dump at /tmp/heapdump.hprof"),
        ("Prevention",
            "Set memory limit to 2Gi · Add readiness probe · "
            "Configure HPA with memory threshold at 80%"),
    ]
    y = 2.75
    for sec_title, sec_body in sections:
        add_text(slide, 0.75, y, 2.5, 0.35,
                 sec_title, 11, bold=True, color=RED)
        add_text(slide, 0.75, y+0.35, 11.6, 1.05,
                 sec_body, 11, color=BLACK, wrap=True)
        add_rect(slide, 0.75, y+1.4, 11.7, 0.04, GRAY40)
        y += 1.5

    add_text(slide, 0.6, 6.5, 12.0, 0.35,
             "Confidence: HIGH  ·  Generated in < 60 seconds  ·  Consistent structure every time",
             11, italic=True, color=GRAY60, align=PP_ALIGN.CENTER)


def slide_08_platform(slide):
    add_content_chrome(slide, "Platform Fit", 8)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "Runs on Technology You Already Own", 22,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 9.0, 0.4,
             "No new vendors, no new licenses, no new infrastructure", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    components = [
        ("OpenShift", "Runs inside the cluster — native RBAC, no VPN required"),
        ("Red Hat OpenShift AI", "llama-3.2-3B via vLLM — self-hosted, air-gap capable"),
        ("WildFly / JBoss EAP", "Operator-managed — watches any WildFly CR"),
        ("Prometheus / AlertManager", "Polls existing monitoring stack — no new exporters"),
        ("JIRA Cloud", "REST API v3 — ADF-formatted rich tickets"),
    ]
    icons = ["●", "●", "●", "●", "●"]
    col_colors = [DARK_PUR, RED, TEAL, GRAY60, DARK_PUR]
    for i, ((comp, desc), col) in enumerate(zip(components, col_colors)):
        y = 2.15 + i * 0.87
        add_rect(slide, 0.6, y, 0.55, 0.68, col)
        add_text(slide, 0.6, y+0.12, 0.55, 0.45,
                 icons[i], 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.2, y, 11.2, 0.68, GRAY20 if i % 2 == 0 else WHITE)
        add_text(slide, 1.35, y+0.06, 3.5, 0.35, comp, 14, bold=True, color=col)
        add_text(slide, 5.0, y+0.1, 7.2, 0.55, desc, 13, color=BLACK)


def slide_09_competitive(slide):
    add_content_chrome(slide, "Differentiation", 9)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "No Generic APM Tool Does This", 22,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 9.0, 0.4,
             "JBoss-specific AI analysis — not alerts, not dashboards, not runbooks", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    features = [
        "Detects JBoss-specific crash patterns (OOMKilled, CrashLoopBackOff, WFLYCTL errors)",
        "AI generates WildFly/EAP-specific root cause analysis — not generic advice",
        "Forces structured JSON output via function-calling — no free-text hallucination",
        "Deduplication prevents ticket storms during sustained incidents",
        "Self-hosted RHOAI model — data never leaves your cluster",
        "Open source — Apache 2.0, fully auditable, no black-box vendor lock-in",
    ]
    check_colors = [RED, RED, RED, DARK_PUR, DARK_PUR, TEAL]
    for i, (feat, cc) in enumerate(zip(features, check_colors)):
        y = 2.15 + i * 0.72
        add_rect(slide, 0.6, y, 0.55, 0.58, cc)
        add_text(slide, 0.6, y+0.1, 0.55, 0.4, "✓", 16,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 1.3, y+0.06, 11.5, 0.55, feat, 14, color=BLACK)


def slide_10_timeline(slide):
    add_content_chrome(slide, "Getting Started", 10)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "Running in Your Environment in Under a Day", 22,
             bold=True, color=BLACK)
    add_text(slide, 0.97, 1.4, 9.0, 0.4,
             "Minimal prerequisites — works alongside your existing OpenShift cluster", 14,
             italic=True, color=GRAY60)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    phases = [
        ("Hour 1", "Install WildFly Operator + deploy WildFlyServer CR", RED),
        ("Hour 2", "Configure credentials (RHOAI token, JIRA API key)", DARK_PUR),
        ("Hour 3", "Build & push container image, apply ConfigMap", TEAL),
        ("Hour 4", "Deploy monitor pod, verify JIRA tickets flowing", GREEN),
        ("Day 2+", "Tune alert filters, dedup window, LLM model", GRAY60),
    ]
    for i, (phase, desc, col) in enumerate(phases):
        y = 2.15 + i * 0.82
        # Phase badge
        add_rect(slide, 0.6, y, 2.0, 0.68, col)
        add_text(slide, 0.6, y+0.14, 2.0, 0.4,
                 phase, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Arrow connector
        add_text(slide, 2.65, y+0.16, 0.4, 0.4, "→", 16, color=GRAY40)
        # Description
        add_rect(slide, 3.1, y, 9.7, 0.68, GRAY20 if i % 2 == 0 else WHITE)
        add_text(slide, 3.25, y+0.14, 9.4, 0.45, desc, 14, color=BLACK)


def slide_11_cta(slide):
    add_content_chrome(slide, "Next Steps", 11)
    add_text(slide, 0.97, 0.7, 11.4, 0.65,
             "Three Ways to Move Forward Today", 22,
             bold=True, color=BLACK)
    add_rect(slide, 0.97, 1.9, 11.4, 0.04, RED)

    options = [
        ("Pilot", "Deploy on one non-production JBoss instance in 4 hours.\nWe provide the container image and all YAML.\nMeasure time-to-first-ticket within the same day.",
         RED),
        ("Assessment", "90-minute architecture review of your current OpenShift + RHOAI setup.\nIdentify quick wins and quantify the ROI for your specific environment.\nNo commitment required.",
         DARK_PUR),
        ("Production\nRollout", "Phased deployment across all JBoss/EAP instances.\nCustom alert filter tuning + JIRA project configuration.\nFull SRE handoff and runbook documentation.",
         TEAL),
    ]
    for i, (title, body, col) in enumerate(options):
        x = 0.6 + i * 4.25
        add_rect(slide, x, 2.1, 4.0, 4.8, col)
        add_rect(slide, x, 2.1, 4.0, 0.08, WHITE)
        add_text(slide, x+0.15, 2.25, 3.7, 0.7,
                 title, 18, bold=True, color=WHITE, wrap=True)
        add_rect(slide, x+0.15, 2.95, 3.7, 0.04, RGBColor(0xFF, 0xFF, 0xFF))
        add_text(slide, x+0.15, 3.1, 3.7, 3.6,
                 body, 13, color=WHITE, wrap=True)

    add_text(slide, 0.6, 6.5, 12.1, 0.35,
             "Source code: github.com/nnsaraiya/jboss-ai-agent-monitor  ·  Apache 2.0",
             11, italic=True, color=GRAY60, align=PP_ALIGN.CENTER)


def slide_12_closing(slide):
    """Thank you / questions — dark background."""
    add_rect(slide, 0, 0, W, H, DARK_PUR)
    add_rect(slide, 0, 0, 0.18, H, RED)
    add_text(slide, 0.5, 1.8, 11.0, 0.6,
             "Thank You", 48, bold=True, color=WHITE)
    add_rect(slide, 0.5, 2.65, 6.0, 0.06, RED)
    add_text(slide, 0.5, 2.9, 11.0, 0.6,
             "Questions?", 36, color=PURPLE)
    add_text(slide, 0.5, 3.8, 8.0, 0.45,
             "Nilay Saraiya  ·  Senior Architect, Red Hat", 16, color=GRAY40)
    add_text(slide, 0.5, 4.35, 8.0, 0.45,
             "nilay.saraiya@redhat.com", 14, color=GRAY40)
    add_text(slide, 0.5, 5.1, 10.0, 0.45,
             "github.com/nnsaraiya/jboss-ai-agent-monitor  ·  Apache 2.0", 13,
             color=GRAY60)
    add_rect(slide, 11.2, 6.4, 1.93, 0.6, RED)
    add_text(slide, 11.25, 6.45, 1.85, 0.5,
             "RED HAT", 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Slide number
    add_text(slide, 0.09, 6.75, 0.8, 0.25, "12", 9, color=GRAY40)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    blank = prs.slide_layouts[6]   # blank layout — no placeholders
    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_solution,
        slide_04_roi,
        slide_05_scale,
        slide_06_oncall,
        slide_07_ticket,
        slide_08_platform,
        slide_09_competitive,
        slide_10_timeline,
        slide_11_cta,
        slide_12_closing,
    ]

    for builder in builders:
        slide = prs.slides.add_slide(blank)
        builder(slide)

    out = "docs/jboss-ai-monitor-deck.pptx"
    prs.save(out)
    print(f"DONE — deck written to {out}")


if __name__ == "__main__":
    main()
