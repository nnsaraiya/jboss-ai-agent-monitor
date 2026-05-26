#!/usr/bin/env python3
"""
generate_deck_v2.py — Rebuild jboss-ai-monitor-deck.pptx (12 slides)
Matches the Red Hat visual design system described in the brief.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────────────────
RED        = RGBColor(0xEE, 0x00, 0x00)
BLUE       = RGBColor(0x3C, 0x8F, 0xF7)
DARK_GRAY  = RGBColor(0x21, 0x21, 0x21)
MID_GRAY   = RGBColor(0x59, 0x59, 0x59)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x3E, 0x8C, 0x35)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_BG    = RGBColor(0x1E, 0x1E, 0x1E)
ALT_GRAY   = RGBColor(0xDD, 0xDD, 0xDD)
ROW_ALT    = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_RED_TINT = RGBColor(0xFF, 0xF0, 0xF0)
FOOTER_BG  = RGBColor(0x21, 0x21, 0x21)

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size, bold, color_rgb,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.text_frame.word_wrap = wrap
    # Handle multi-line text split by \n
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color_rgb
        try:
            run.font.name = 'Calibri'
        except Exception:
            pass
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size, color_rgb, bullet_char="▶  "):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color_rgb
        try:
            run.font.name = 'Calibri'
        except Exception:
            pass
    return txBox


def add_red_bar_top(slide):
    """Thin red bar across the top of the slide."""
    add_rect(slide, 0, 0, 20, 0.5, RED)


def cell(slide, label, x, y, w, h, fill_rgb, text_rgb, font_size, bold=False,
         align=PP_ALIGN.LEFT, italic=False):
    """Draw a table cell (rect + text)."""
    add_rect(slide, x, y, w, h, fill_rgb)
    pad = 0.1
    add_text_box(slide, label, x + pad, y + 0.05, w - 2 * pad, h - 0.1,
                 font_size, bold, text_rgb, align=align, italic=italic)


# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(20)
prs.slide_height = Inches(11.25)

BLANK = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (dark background)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full dark background
add_rect(slide, 0, 0, 20, 11.25, NEAR_BLACK)
# Left red accent bar
add_rect(slide, 0, 0, 0.5, 11.25, RED)
# Main title
add_text_box(slide,
             "Autonomous JBoss Operations\nwith Red Hat OpenShift AI",
             1.2, 2.0, 17, 3.5,
             60, True, WHITE)
# Subtitle
add_text_box(slide,
             "Reducing MTTR and operational cost with an AI-powered monitoring agent",
             1.2, 5.8, 14, 1.5,
             28, False, LIGHT_GRAY)
# Bottom small text
add_text_box(slide, "Red Hat — Confidential",
             14, 10.5, 5.5, 0.5,
             16, False, MID_GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide,
             "The Status Quo: Manual, Reactive, Expensive",
             1.0, 0.7, 18, 2.5,
             60, True, RED)
# Thin red rule
add_rect(slide, 1.0, 3.0, 18, 0.06, RED)
# Bullet list
add_bullet_list(slide,
    ["JBoss incidents go undetected until users complain",
     "Senior engineers spend 2+ hours per incident on root cause",
     "JIRA tickets created manually — late, incomplete, inconsistent",
     "On-call rotations are expensive and cause burnout",
     "The same incidents recur — fixes are never documented"],
    1.5, 3.3, 17, 7.0,
    28, DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution (4 pillars)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "Detect. Analyse. Document. Automatically.",
             1.0, 0.7, 18, 2.0,
             60, True, RED)

pillars = [
    (0.5,  "\U0001f50d", "Detect",      "Monitors pod state, logs, alerts and health every 60 seconds"),
    (5.0,  "\U0001f9e0", "Analyse",     "RHOAI-hosted LLM generates root cause and resolution steps"),
    (9.5,  "\U0001f3ab", "Document",    "Structured JIRA ticket created automatically with full context"),
    (14.0, "\U0001f515", "Deduplicate", "Suppresses repeat tickets within a configurable time window"),
]

for bx, emoji, title, body in pillars:
    bw = 4.3
    by = 3.2
    bh = 7.0
    # Main box
    add_rect(slide, bx, by, bw, bh, LIGHT_GRAY)
    # Red top strip
    add_rect(slide, bx, by, bw, 0.15, RED)
    # Emoji
    add_text_box(slide, emoji, bx + 0.1, 3.6, bw - 0.2, 1.2,
                 44, False, DARK_GRAY, align=PP_ALIGN.CENTER)
    # Box title
    add_text_box(slide, title, bx + 0.1, 5.0, bw - 0.2, 0.7,
                 24, True, RED, align=PP_ALIGN.CENTER)
    # Body text
    add_text_box(slide, body, bx + 0.2, 5.8, bw - 0.4, 4.2,
                 22, False, DARK_GRAY, align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Time Savings (table)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide,
             "2.5 Hours of Senior Engineer Time → 15 Minutes of Review",
             1.0, 0.7, 18, 1.8,
             52, True, RED)

# Table column config: (x, width)
cols = [(1.0, 5.0), (6.2, 4.0), (10.4, 4.0), (14.6, 4.8)]

# Header row
header_labels = ["Step", "Before", "After", "Time Saved"]
for (cx, cw), lbl in zip(cols, header_labels):
    cell(slide, lbl, cx, 2.8, cw, 0.7, RED, WHITE, 22, bold=True, align=PP_ALIGN.CENTER)

# Data rows
data_rows = [
    ("Issue detection",      "5–30 min",   "60 seconds",     "~29 min saved"),
    ("Root cause analysis",  "60–120 min", "AI: 14 seconds", "~119 min saved"),
    ("JIRA ticket creation", "10–15 min",  "Automated",      "~13 min saved"),
]
row_fills = [WHITE, RGBColor(0xF5, 0xF5, 0xF5), WHITE]
row_ys    = [3.6, 4.5, 5.4]

for row_y, row_fill, row_data in zip(row_ys, row_fills, data_rows):
    for (cx, cw), txt in zip(cols, row_data):
        cell(slide, txt, cx, row_y, cw, 0.8, row_fill, DARK_GRAY, 22,
             align=PP_ALIGN.CENTER)

# Footer row
footer_data = ["TOTAL", "~2.5 hours", "~15 min review", "~2 hours saved"]
for (cx, cw), txt in zip(cols, footer_data):
    cell(slide, txt, cx, 6.3, cw, 0.9, FOOTER_BG, WHITE, 24, bold=True,
         align=PP_ALIGN.CENTER)

# Callout box
add_rect(slide, 3.0, 7.6, 14, 1.2, BLUE)
add_text_box(slide,
             "At $100/hr blended SRE rate → $225 saved per incident",
             3.0, 7.65, 14, 1.1,
             26, True, WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROI by Size
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "The Value Scales With Your Environment",
             1.0, 0.7, 18, 2.0,
             60, True, RED)

# Box 1 — Small
add_rect(slide, 0.5, 3.2, 6.0, 6.5, LIGHT_GRAY)
add_text_box(slide, "Small — 1 env · 10 incidents/mo",
             0.6, 3.8, 5.8, 0.9, 24, False, DARK_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, "$27,000/yr",
             0.6, 5.5, 5.8, 1.2, 48, True, RED, align=PP_ALIGN.CENTER)
add_text_box(slide, "annual savings",
             0.6, 7.2, 5.8, 0.8, 20, False, MID_GRAY, align=PP_ALIGN.CENTER)

# Box 2 — Medium
add_rect(slide, 7.0, 3.2, 6.0, 6.5, ALT_GRAY)
add_text_box(slide, "Medium — 3 envs · 30 incidents/mo",
             7.1, 3.8, 5.8, 0.9, 24, False, DARK_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, "$81,000/yr",
             7.1, 5.5, 5.8, 1.2, 48, True, RED, align=PP_ALIGN.CENTER)
add_text_box(slide, "annual savings",
             7.1, 7.2, 5.8, 0.8, 20, False, MID_GRAY, align=PP_ALIGN.CENTER)

# Box 3 — Large (red background, white text)
add_rect(slide, 13.5, 3.2, 6.0, 6.5, RED)
add_text_box(slide, "Large — 10 envs · 100 incidents/mo",
             13.6, 3.8, 5.8, 0.9, 24, False, WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "$270,000/yr",
             13.6, 5.5, 5.8, 1.2, 48, True, WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "annual savings",
             13.6, 7.2, 5.8, 0.8, 20, False, WHITE, align=PP_ALIGN.CENTER)

# Footer note
add_text_box(slide,
             "Conservative estimate — direct engineer time only. "
             "Excludes on-call premiums, prevented downtime, and revenue protection.",
             1.0, 10.3, 18, 0.7,
             18, False, MID_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — On-Call Value
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "24/7 Coverage — Without 3am Pages",
             1.0, 0.7, 18, 2.0,
             60, True, RED)

# Dark quote box
add_rect(slide, 1.0, 3.0, 18, 3.5, FOOTER_BG)
add_text_box(slide,
             "“The first P1 incident it catches at 3am on a Sunday — without waking anyone "
             "— pays for the entire implementation.”",
             1.2, 3.2, 17.6, 3.1,
             30, False, WHITE, align=PP_ALIGN.CENTER, italic=True)

# 3 stat boxes
stat_boxes = [
    (0.5,  "10–25%", "On-call salary premium eliminated"),
    (7.3,  "3am",         "Agent works nights, weekends, holidays"),
    (14.0, "0 min",       "Time to ticket after issue detection"),
]
for bx, big_label, desc in stat_boxes:
    add_rect(slide, bx, 7.0, 5.5, 3.5, LIGHT_GRAY)
    add_text_box(slide, big_label,
                 bx + 0.1, 7.4, 5.3, 1.2,
                 40, True, RED, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc,
                 bx + 0.2, 8.7, 5.1, 1.6,
                 22, False, DARK_GRAY, align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Sample JIRA Ticket
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide,
             "Every Incident Gets a Structured, Actionable Ticket",
             1.0, 0.7, 14, 2.0,
             52, True, RED)

# Dark ticket box
add_rect(slide, 1.0, 3.0, 13.5, 7.8, DARK_BG)

ticket_text = (
    "\U0001f534  OOMKilled — jboss-instance-0  |  CRITICAL  |  jboss-production\n"
    "\n"
    "\U0001f9e0  Root Cause  (AI — Confidence: HIGH)\n"
    "JVM heap exhausted. -Xmx setting too low for workload.\n"
    "Container memory limit: 1536Mi.  Current heap: 256Mi.\n"
    "\n"
    "\U0001f527  Resolution Steps\n"
    "  1. Increase -Xmx to 1024m in standalone.xml\n"
    "  2. Set OpenShift memory limit to 1536Mi\n"
    "  3. Enable -XX:+HeapDumpOnOutOfMemoryError\n"
    "\n"
    "\U0001f6e1️  Prevention Tips\n"
    "  •  Set JVM heap to 75% of container memory limit\n"
    "  •  Add readiness probe on /health/ready\n"
    "\n"
    "\U0001f4da  References\n"
    "  •  Red Hat KBase: Tuning JVM for EAP on OpenShift"
)
add_text_box(slide, ticket_text,
             1.5, 3.3, 12.5, 7.2,
             20, False, WHITE, wrap=True)

# Right callout box (green)
add_rect(slide, 15.0, 4.5, 4.5, 3.0, GREEN)
add_text_box(slide,
             "Created in\n14 seconds\nNo human\ninvolved",
             15.0, 4.6, 4.5, 2.8,
             26, True, WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Red Hat Platform
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "Runs on Technology You Already Own",
             1.0, 0.7, 18, 2.0,
             60, True, RED)

# Left table — header
platform_cols = [(1.0, 5.5), (6.7, 4.8)]
platform_header = ["Component", "Red Hat Product"]
for (cx, cw), lbl in zip(platform_cols, platform_header):
    cell(slide, lbl, cx, 3.0, cw, 1.0, RED, WHITE, 22, bold=True,
         align=PP_ALIGN.CENTER)

platform_rows = [
    ("JBoss runtime",      "Red Hat JBoss EAP / WildFly Operator"),
    ("Container platform", "Red Hat OpenShift"),
    ("AI inference",       "Red Hat OpenShift AI (RHOAI)"),
    ("Model serving",      "KServe + vLLM on GPU"),
    ("Container image",    "UBI9 (Universal Base Image)"),
]
row_fills_p = [WHITE, LIGHT_GRAY, WHITE, LIGHT_GRAY, WHITE]
for i, (row_fill, (comp, prod)) in enumerate(zip(row_fills_p, platform_rows)):
    ry = 4.1 + i * 1.1
    cell(slide, comp, 1.0, ry, 5.5, 1.0, row_fill, DARK_GRAY, 22)
    cell(slide, prod, 6.7, ry, 4.8, 1.0, row_fill, DARK_GRAY, 22)

# Right callout boxes (blue)
callouts = [
    (3.3, "✓  Activates your RHOAI investment on Day 1"),
    (5.3, "✓  No new vendor relationship"),
    (7.3, "✓  Customer data never leaves the cluster"),
]
for cy, txt in callouts:
    add_rect(slide, 12.0, cy, 7.5, 1.6, BLUE)
    add_text_box(slide, txt, 12.2, cy + 0.2, 7.1, 1.2,
                 24, False, WHITE, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Competitive Differentiation
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "No Generic APM Tool Does This",
             1.0, 0.7, 18, 2.0,
             55, True, RED)

# Header row — 4 columns
comp_cols = [
    (0.5,  7.0),
    (7.7,  3.5),
    (11.4, 3.5),
    (15.1, 4.5),
]
comp_header_labels = ["Capability", "Generic APM", "Cloud AI Services", "This Solution"]
for (cx, cw), lbl in zip(comp_cols, comp_header_labels):
    cell(slide, lbl, cx, 3.0, cw, 0.8, RED, WHITE, 22, bold=True, align=PP_ALIGN.CENTER)

comp_rows = [
    ("JBoss-specific root cause",               "✗",  "Partial",  "✓"),
    ("AI resolution steps",                     "✗",  "✗",   "✓"),
    ("On-premises LLM (no data leaves cluster)", "✗", "✗",   "✓"),
    ("Structured JIRA tickets",                 "✗",  "✗",   "✓"),
    ("No per-API-call cost",                    "—",  "✗",   "✓"),
    ("OpenShift-native",                        "✗",  "✗",   "✓"),
]
row_fills_c = [WHITE, ROW_ALT, WHITE, ROW_ALT, WHITE, ROW_ALT]

for i, (row_fill, row_data) in enumerate(zip(row_fills_c, comp_rows)):
    ry = 3.9 + i * 0.9
    cap, apm_val, cloud_val, ours = row_data

    # Capability cell
    cell(slide, cap,       0.5,  ry, 7.0, 0.9, row_fill, DARK_GRAY, 21)
    # APM cell
    apm_color = MID_GRAY if apm_val == "✗" else DARK_GRAY
    cell(slide, apm_val,   7.7,  ry, 3.5, 0.9, row_fill, apm_color, 21,
         align=PP_ALIGN.CENTER)
    # Cloud cell
    cloud_color = MID_GRAY if cloud_val == "✗" else DARK_GRAY
    cell(slide, cloud_val, 11.4, ry, 3.5, 0.9, row_fill, cloud_color, 21,
         align=PP_ALIGN.CENTER)
    # "This Solution" cell — light red tint, green checkmark
    our_color = GREEN if ours == "✓" else DARK_GRAY
    our_bold  = (ours == "✓")
    cell(slide, ours,      15.1, ry, 4.5, 0.9, LIGHT_RED_TINT, our_color, 21,
         bold=our_bold, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Time to Value
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "Running in Your Environment in Under a Day",
             1.0, 0.7, 18, 2.0,
             60, True, RED)

phases = [
    (0.5,  LIGHT_GRAY, BLUE,  RED,       DARK_GRAY,
     "PHASE 1 — 2 to 4 hours", "Deploy",
     "OpenShift manifests, RHOAI endpoint, JIRA configuration"),
    (7.3,  LIGHT_GRAY, BLUE,  RED,       DARK_GRAY,
     "PHASE 2 — 1 to 2 hours", "Validate",
     "Inject test errors, confirm ticket creation, tune patterns"),
    (14.1, RED,        WHITE, WHITE,     WHITE,
     "PHASE 3 — Ongoing",       "Expand",
     "Add namespaces, extend to Quarkus and Spring Boot"),
]

for bx, bg, lbl_col, title_col, body_col, phase_lbl, phase_title, body in phases:
    add_rect(slide, bx, 3.2, 5.5, 6.8, bg)
    add_text_box(slide, phase_lbl,
                 bx + 0.15, 3.5, 5.2, 0.7, 22, True, lbl_col)
    add_text_box(slide, phase_title,
                 bx + 0.15, 4.4, 5.2, 0.9, 30, True, title_col)
    add_text_box(slide, body,
                 bx + 0.15, 5.4, 5.2, 4.3, 22, False, body_col, wrap=True)

# Arrows between boxes
add_text_box(slide, "→", 6.3, 6.0, 0.9, 0.9, 32, True, RED, align=PP_ALIGN.CENTER)
add_text_box(slide, "→", 13.1, 6.0, 0.9, 0.9, 32, True, RED, align=PP_ALIGN.CENTER)

# Footer
add_text_box(slide,
             "No professional services engagement required for initial deployment. "
             "Customer owns and operates it.",
             1.0, 10.3, 18, 0.7,
             20, False, MID_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Call to Action
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_red_bar_top(slide)
add_text_box(slide, "Three Ways to Move Forward Today",
             1.0, 0.7, 18, 2.0,
             60, True, RED)

options = [
    (3.0, LIGHT_GRAY, BLUE,  DARK_GRAY,
     "Option 1 — Proof of Value (1 week)",
     "Deploy in non-production. Run for one week. Measure MTTR delta vs baseline."),
    (5.4, LIGHT_GRAY, BLUE,  DARK_GRAY,
     "Option 2 — Architecture Workshop (half day)",
     "Red Hat architect maps the solution to your JBoss topology and RHOAI setup."),
    (7.8, RED,        WHITE, WHITE,
     "Option 3 — Production Deployment",
     "Scoped engagement: deploy, configure, and hand over with runbook and training."),
]

for oy, bg, lbl_col, body_col, lbl, body in options:
    add_rect(slide, 1.5, oy, 17, 2.0, bg)
    add_text_box(slide, lbl,  1.7, oy + 0.15, 16.6, 0.7, 22, True,  lbl_col)
    add_text_box(slide, body, 1.7, oy + 0.9,  16.6, 0.9, 22, False, body_col, wrap=True)

# Footer
add_text_box(slide,
             "All source code is open and customer-owned.",
             1.5, 10.2, 17, 0.6,
             20, False, MID_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full dark background
add_rect(slide, 0, 0, 20, 11.25, NEAR_BLACK)
# Red left bar
add_rect(slide, 0, 0, 0.5, 11.25, RED)
# Big white "Questions?"
add_text_box(slide, "Questions?",
             1.5, 2.5, 16, 3.0,
             80, True, WHITE)
# Subtitle
add_text_box(slide, "Let us talk about your JBoss environment.",
             1.5, 5.8, 14, 1.5,
             32, False, LIGHT_GRAY)
# Contact placeholder
add_text_box(slide,
             "[Name] · [email] · Red Hat Senior Architect",
             1.5, 9.5, 14, 0.8,
             22, False, MID_GRAY)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "/home/user/jboss-ai-agent-monitor/docs/jboss-ai-monitor-deck.pptx"
prs.save(out_path)
print("DONE — deck written to docs/jboss-ai-monitor-deck.pptx")
