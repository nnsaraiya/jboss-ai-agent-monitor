from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colors ──────────────────────────────────────────────────
RH_RED   = RGBColor(0xEE, 0x00, 0x00)
DARK_GR  = RGBColor(0x21, 0x21, 0x21)
MID_GR   = RGBColor(0x44, 0x44, 0x44)
LIGHT_GR = RGBColor(0xF5, 0xF5, 0xF5)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x3E, 0x8C, 0x35)
SLIDE_BG = RGBColor(0x1A, 0x1A, 0x1A)

# ── Presentation setup ────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank


# ── Helper: add filled rectangle ──────────────────────────────────
def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()   # no border
    return shape


# ── Helper: add text box ──────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_size, bold, color_rgb,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    # clear the default empty paragraph
    tf.paragraphs[0].text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color_rgb
    run.font.italic = italic
    run.font.name  = font_name
    return tf


# ── Helper: add text paragraph to existing tf ─────────────────────
def add_para(tf, text, font_size, bold, color_rgb,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color_rgb
    run.font.italic = italic
    return p


# ── Helper: full-width red top bar + title ────────────────────────
def add_top_bar(slide, title_text):
    add_rect(slide, 0, 0, 13.33, 0.7, RH_RED)
    add_text(slide, title_text, 0.25, 0.08, 12.8, 0.6,
             font_size=22, bold=True, color_rgb=WHITE, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, SLIDE_BG)          # dark bg
add_rect(s1, 0, 0, 0.3,   7.5, RH_RED)            # red left bar

add_text(s1,
         "Autonomous JBoss Operations\nwith Red Hat OpenShift AI",
         0.6, 2.4, 10.5, 1.8,
         font_size=40, bold=True, color_rgb=WHITE)

add_text(s1,
         "Reducing MTTR and operational cost\nwith an AI-powered monitoring agent",
         0.6, 4.4, 10.5, 1.2,
         font_size=20, bold=False, color_rgb=LIGHT_GR)

add_text(s1, "Red Hat — Confidential",
         10.5, 7.05, 2.6, 0.35,
         font_size=12, bold=False, color_rgb=MID_GR, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s2, "The Status Quo: Manual, Reactive, Expensive")

bullets = [
    "JBoss incidents go undetected until users complain",
    "Senior engineers spend 2+ hours per incident on root cause",
    "JIRA tickets are created manually — late, incomplete, inconsistent",
    "On-call rotations are expensive and cause burnout",
    "The same incidents recur — fixes never get documented",
]

for i, b in enumerate(bullets):
    add_text(s2, f"▶  {b}",
             0.5, 1.0 + i * 1.1, 12.5, 0.9,
             font_size=18, bold=False, color_rgb=DARK_GR)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution (4 pillars)
# ═══════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s3, "The JBoss AI Monitor: Detect. Analyse. Document. Automatically.")

pillars = [
    ("🔍", "Detect",      "Monitors pod state, logs, alerts and health every 60 seconds"),
    ("🧠", "Analyse",     "RHOAI-hosted LLM generates root cause + resolution steps"),
    ("🎫", "Document",    "Structured JIRA ticket created automatically with full context"),
    ("🔕", "Deduplicate", "Suppresses repeat tickets within a configurable time window"),
]

box_w = 3.0
gap   = 0.16
start = 0.22

for i, (icon, title, body) in enumerate(pillars):
    bx = start + i * (box_w + gap)
    by = 1.0
    bh = 5.8

    box = add_rect(s3, bx, by, box_w, bh, LIGHT_GR)

    # red top accent line
    add_rect(s3, bx, by, box_w, 0.07, RH_RED)

    add_text(s3, icon, bx + 0.1, by + 0.2, box_w - 0.2, 0.75,
             font_size=40, bold=False, color_rgb=DARK_GR, align=PP_ALIGN.CENTER)

    add_text(s3, title, bx + 0.1, by + 1.05, box_w - 0.2, 0.45,
             font_size=16, bold=True, color_rgb=RH_RED, align=PP_ALIGN.CENTER)

    add_text(s3, body, bx + 0.15, by + 1.6, box_w - 0.3, 2.0,
             font_size=13, bold=False, color_rgb=MID_GR, align=PP_ALIGN.CENTER, wrap=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Time Savings (table)
# ═══════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s4, "2.5 Hours of Senior Engineer Time → 15 Minutes of Review")

cols  = [3.6, 2.5, 2.8, 2.8]
hdrs  = ["Step", "Before", "After", "Time Saved"]
rows  = [
    ("Issue detection",         "5–30 min",      "60 seconds",     "~29 min saved"),
    ("Root cause analysis",     "60–120 min",     "AI: 14 seconds", "~119 min saved"),
    ("JIRA ticket creation",    "10–15 min",      "Automated",      "~13 min saved"),
]
footer = ("TOTAL", "~2.5 hours", "~15 min review", "~2 hours saved")

row_h  = 0.65
tbl_y  = 0.85
tbl_x  = 0.25

# header row
cx = tbl_x
for j, (h, w) in enumerate(zip(hdrs, cols)):
    add_rect(s4, cx, tbl_y, w, row_h, RH_RED)
    add_text(s4, h, cx + 0.05, tbl_y + 0.1, w - 0.1, row_h - 0.1,
             font_size=14, bold=True, color_rgb=WHITE, align=PP_ALIGN.CENTER)
    cx += w

# data rows
for ri, row in enumerate(rows):
    cx = tbl_x
    bg = WHITE if ri % 2 == 0 else LIGHT_GR
    for j, (cell, w) in enumerate(zip(row, cols)):
        add_rect(s4, cx, tbl_y + row_h * (ri + 1), w, row_h, bg)
        add_text(s4, cell, cx + 0.05,
                 tbl_y + row_h * (ri + 1) + 0.1, w - 0.1, row_h - 0.1,
                 font_size=13, bold=False, color_rgb=DARK_GR, align=PP_ALIGN.CENTER)
        cx += w

# footer row
cx = tbl_x
fy = tbl_y + row_h * (len(rows) + 1)
for j, (cell, w) in enumerate(zip(footer, cols)):
    add_rect(s4, cx, fy, w, row_h, DARK_GR)
    add_text(s4, cell, cx + 0.05, fy + 0.1, w - 0.1, row_h - 0.1,
             font_size=14, bold=True, color_rgb=WHITE, align=PP_ALIGN.CENTER)
    cx += w

# callout box below table
call_y = fy + row_h + 0.2
add_rect(s4, tbl_x, call_y, sum(cols), 0.7, RGBColor(0xFF, 0xDD, 0xDD))
add_text(s4,
         "At $100/hr blended SRE rate → $225 saved per incident",
         tbl_x + 0.15, call_y + 0.12, sum(cols) - 0.3, 0.5,
         font_size=14, bold=True, color_rgb=RH_RED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — ROI by Size
# ═══════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s5, "The Value Scales With Your Environment")

boxes = [
    (LIGHT_GR,           DARK_GR, "Small\n1 env · 10 incidents/mo",     "$27,000/yr saved"),
    (RGBColor(0xE8,0xE8,0xE8), DARK_GR, "Medium\n3 envs · 30 incidents/mo",  "$81,000/yr saved"),
    (RH_RED,             WHITE,   "Large\n10 envs · 100 incidents/mo",  "$270,000/yr saved"),
]

bw = 3.8
bh = 4.8
by = 1.0
bx_start = 0.5

for i, (bg, fg, desc, amount) in enumerate(boxes):
    bx = bx_start + i * (bw + 0.55)
    add_rect(s5, bx, by, bw, bh, bg)
    add_text(s5, desc, bx + 0.2, by + 0.5, bw - 0.4, 1.2,
             font_size=14, bold=False, color_rgb=fg, align=PP_ALIGN.CENTER)
    add_text(s5, amount, bx + 0.2, by + 2.2, bw - 0.4, 1.0,
             font_size=32, bold=True, color_rgb=fg, align=PP_ALIGN.CENTER)

add_text(s5,
         "Conservative estimate — direct engineer time only. Does not include on-call premiums, prevented downtime, or revenue protection.",
         0.4, 6.35, 12.5, 0.7,
         font_size=12, bold=False, color_rgb=MID_GR, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — On-Call Value
# ═══════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s6, "24/7 Coverage — Without 3am Pages")

# large quote box
add_rect(s6, 0.5, 0.9, 12.33, 2.8, DARK_GR)
add_text(s6,
         '"The first P1 incident it catches at 3am on a Sunday\n— without waking anyone —\npays for the entire implementation."',
         0.8, 1.0, 11.8, 2.6,
         font_size=22, bold=False, color_rgb=WHITE,
         align=PP_ALIGN.CENTER, italic=True)

# 3 stat boxes
stats = [
    ("10–25%",  "On-call salary premium eliminated"),
    ("3am",     "Agent works nights, weekends, holidays"),
    ("0 min",   "Time to ticket after detection"),
]
sw = 3.8
sx_start = 0.55
sy = 4.0

for i, (label, desc) in enumerate(stats):
    sx = sx_start + i * (sw + 0.5)
    add_rect(s6, sx, sy, sw, 2.8, LIGHT_GR)
    add_text(s6, label, sx + 0.1, sy + 0.3, sw - 0.2, 1.0,
             font_size=36, bold=True, color_rgb=RH_RED, align=PP_ALIGN.CENTER)
    add_text(s6, desc, sx + 0.1, sy + 1.4, sw - 0.2, 1.2,
             font_size=14, bold=False, color_rgb=DARK_GR, align=PP_ALIGN.CENTER, wrap=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Sample JIRA Ticket
# ═══════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s7, "Every Incident Gets a Structured, Actionable Ticket")

# dark code box
add_rect(s7, 0.3, 0.85, 9.8, 6.4, RGBColor(0x1E, 0x1E, 0x1E))

ticket_lines = [
    (WHITE,   True,  18, "🔴 OOMKilled — jboss-instance-0  |  CRITICAL  |  jboss-production"),
    (WHITE,   False, 10, ""),
    (GREEN,   True,  12, "🧠 Root Cause (AI — Confidence: HIGH)"),
    (WHITE,   False, 11, "JVM heap exhausted. -Xmx setting too low for workload."),
    (WHITE,   False, 11, "Container memory limit: 1536Mi. Current heap: 256Mi."),
    (WHITE,   False, 10, ""),
    (GREEN,   True,  12, "🔧 Resolution Steps"),
    (WHITE,   False, 11, "1.  Increase -Xmx to 1024m in standalone.xml"),
    (WHITE,   False, 11, "2.  Set OpenShift memory limit to 1536Mi"),
    (WHITE,   False, 11, "3.  Enable -XX:+HeapDumpOnOutOfMemoryError"),
    (WHITE,   False, 10, ""),
    (GREEN,   True,  12, "🛡️ Prevention Tips"),
    (WHITE,   False, 11, "• Set JVM heap to 75% of container memory limit"),
    (WHITE,   False, 11, "• Add readiness probe on /health/ready"),
    (WHITE,   False, 10, ""),
    (GREEN,   True,  12, "📚 References"),
    (WHITE,   False, 11, "• Red Hat KBase: Tuning JVM for EAP on OpenShift"),
]

# draw each line as separate text box
for li, (col, bld, fsz, txt) in enumerate(ticket_lines):
    add_text(s7, txt,
             0.45, 0.98 + li * 0.33, 9.5, 0.35,
             font_size=fsz, bold=bld, color_rgb=col,
             font_name="Courier New")

# right callout
add_rect(s7, 10.3, 2.5, 2.7, 1.4, RGBColor(0xE8, 0xF5, 0xE8))
add_text(s7, "Created in 14 seconds\nNo human involved",
         10.35, 2.6, 2.6, 1.2,
         font_size=14, bold=True, color_rgb=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Red Hat Platform
# ═══════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s8, "Runs on Technology You Already Own")

# Left column table
comp_rows = [
    ("Component",          "Red Hat Product"),
    ("JBoss runtime",      "JBoss EAP / WildFly Operator"),
    ("Container platform", "Red Hat OpenShift"),
    ("AI inference",       "Red Hat OpenShift AI (RHOAI)"),
    ("Model serving",      "KServe + vLLM on GPU"),
    ("Container image",    "UBI9 (Universal Base Image)"),
]

col1_w = 3.2
col2_w = 4.5
rh_  = 0.65
tblx = 0.35
tbly = 0.85

for ri, (c1, c2) in enumerate(comp_rows):
    bg = RH_RED if ri == 0 else (LIGHT_GR if ri % 2 == 1 else WHITE)
    fg = WHITE if ri == 0 else DARK_GR
    add_rect(s8, tblx,            tbly + ri * rh_, col1_w, rh_, bg)
    add_rect(s8, tblx + col1_w,   tbly + ri * rh_, col2_w, rh_, bg)
    add_text(s8, c1, tblx + 0.08, tbly + ri * rh_ + 0.12, col1_w - 0.1, rh_ - 0.1,
             font_size=13, bold=(ri == 0), color_rgb=fg)
    add_text(s8, c2, tblx + col1_w + 0.08, tbly + ri * rh_ + 0.12, col2_w - 0.1, rh_ - 0.1,
             font_size=13, bold=(ri == 0), color_rgb=fg)

# Right column callouts
callouts = [
    "✓  Activates your RHOAI investment",
    "✓  No new vendor relationship",
    "✓  Data never leaves your cluster",
]
rx = 8.3
ry_start = 1.1
for i, txt in enumerate(callouts):
    ry = ry_start + i * 1.55
    add_rect(s8, rx, ry, 4.7, 1.2, LIGHT_GR)
    add_rect(s8, rx, ry, 0.08, 1.2, RH_RED)
    add_text(s8, txt, rx + 0.2, ry + 0.2, 4.4, 0.85,
             font_size=15, bold=True, color_rgb=DARK_GR)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Competitive Differentiation
# ═══════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s9, "No Generic APM Tool Does This")

cap_rows = [
    ("Capability",                    "Generic APM", "Cloud AI",  "This Solution"),
    ("JBoss-specific root cause",      "✗",           "Partial",   "✓"),
    ("AI resolution steps",            "✗",           "✗",         "✓"),
    ("On-premises LLM",                "✗",           "✗",         "✓"),
    ("Data stays in cluster",          "✗",           "✗",         "✓"),
    ("Structured JIRA tickets",        "✗",           "✗",         "✓"),
    ("No per-API-call cost",           "—",           "✗",         "✓"),
]

cw = [4.8, 2.4, 2.4, 2.6]
crh = 0.74
ctx = 0.3
cty = 0.82

for ri, row in enumerate(cap_rows):
    cx2 = ctx
    for ci, (cell, w) in enumerate(zip(row, cw)):
        is_hdr = (ri == 0)
        is_sol = (ci == 3)
        if is_hdr and is_sol:
            bg = RH_RED; fg = WHITE
        elif is_hdr:
            bg = DARK_GR; fg = WHITE
        elif is_sol:
            bg = RGBColor(0xFF, 0xEE, 0xEE); fg = GREEN if cell == "✓" else DARK_GR
        elif cell == "✗":
            bg = WHITE; fg = MID_GR
        else:
            bg = WHITE; fg = DARK_GR

        add_rect(s9, cx2, cty + ri * crh, w, crh, bg)
        fsz = 13 if not is_hdr else 14
        add_text(s9, cell, cx2 + 0.1, cty + ri * crh + 0.1, w - 0.15, crh - 0.1,
                 font_size=fsz, bold=is_hdr, color_rgb=fg, align=PP_ALIGN.CENTER)
        cx2 += w


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Time to Value
# ═══════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s10, "Running in Your Environment in Under a Day")

phases = [
    ("Phase 1\n2–4 hrs",  "Deploy",
     "OpenShift manifests, RHOAI endpoint, JIRA config"),
    ("Phase 2\n1–2 hrs",  "Validate",
     "Inject test errors, confirm ticket creation, tune patterns"),
    ("Phase 3\nOngoing",  "Expand",
     "Add namespaces, extend to Quarkus / Spring Boot"),
]

pw = 3.6
ph = 4.5
py = 1.1
px_start = 0.5

for i, (phase_lbl, phase_title, phase_body) in enumerate(phases):
    px = px_start + i * (pw + 0.65)
    add_rect(s10, px, py, pw, ph, LIGHT_GR)
    add_rect(s10, px, py, pw, 0.08, RH_RED)

    add_text(s10, phase_lbl, px + 0.1, py + 0.2, pw - 0.2, 0.8,
             font_size=13, bold=True, color_rgb=RH_RED, align=PP_ALIGN.CENTER)
    add_text(s10, phase_title, px + 0.1, py + 1.15, pw - 0.2, 0.6,
             font_size=20, bold=True, color_rgb=DARK_GR, align=PP_ALIGN.CENTER)
    add_text(s10, phase_body, px + 0.15, py + 1.9, pw - 0.3, 2.4,
             font_size=13, bold=False, color_rgb=MID_GR, align=PP_ALIGN.CENTER, wrap=True)

    # arrow between phases (except last)
    if i < len(phases) - 1:
        ax = px + pw + 0.12
        ay = py + ph / 2 - 0.2
        add_text(s10, "→", ax, ay, 0.4, 0.4,
                 font_size=24, bold=True, color_rgb=RH_RED, align=PP_ALIGN.CENTER)

add_text(s10,
         "No professional services engagement required for initial deployment.",
         0.4, 6.1, 12.5, 0.5,
         font_size=13, bold=False, color_rgb=MID_GR, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Call to Action
# ═══════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_rect(s11, 0, 0, 13.33, 7.5, WHITE)
add_top_bar(s11, "Three Ways to Move Forward Today")

options = [
    ("Option 1 — Proof of Value (1 week)",
     "Deploy in non-production. Run for one week. Measure MTTR delta vs baseline."),
    ("Option 2 — Architecture Workshop (half day)",
     "Red Hat architect maps solution to your JBoss topology and RHOAI setup."),
    ("Option 3 — Production Deployment",
     "Scoped engagement: deploy, configure, hand over with runbook and training."),
]

ow = 3.8
oh = 4.0
oy = 1.0
ox_start = 0.4

for i, (opt_title, opt_body) in enumerate(options):
    ox = ox_start + i * (ow + 0.55)
    add_rect(s11, ox, oy, ow, oh, LIGHT_GR)
    add_rect(s11, ox, oy, ow, 0.07, RH_RED)

    add_text(s11, opt_title, ox + 0.15, oy + 0.2, ow - 0.3, 0.8,
             font_size=14, bold=True, color_rgb=RH_RED, wrap=True)
    add_text(s11, opt_body, ox + 0.15, oy + 1.15, ow - 0.3, 2.6,
             font_size=13, bold=False, color_rgb=DARK_GR, wrap=True)

add_text(s11, "All source code is open and customer-owned.",
         0.4, 6.0, 12.5, 0.5,
         font_size=13, bold=True, color_rgb=DARK_GR, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You
# ═══════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_rect(s12, 0, 0, 13.33, 7.5, SLIDE_BG)
add_rect(s12, 0, 0, 0.3, 7.5, RH_RED)

add_text(s12, "Questions?",
         0.6, 2.5, 10.5, 1.4,
         font_size=52, bold=True, color_rgb=WHITE)

add_text(s12, "Let's talk about your JBoss environment.",
         0.6, 4.1, 10.5, 0.9,
         font_size=22, bold=False, color_rgb=LIGHT_GR)

add_text(s12, "[Name]  ·  [email]  ·  Red Hat Senior Architect",
         0.6, 6.7, 12.5, 0.55,
         font_size=14, bold=False, color_rgb=MID_GR)


# ── Save ──────────────────────────────────────────────────────────
out = "/home/user/jboss-ai-agent-monitor/docs/jboss-ai-monitor-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
