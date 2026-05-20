from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────────
RED       = RGBColor(0xEE, 0x00, 0x00)
DARK_GRAY = RGBColor(0x21, 0x21, 0x21)
MID_GRAY  = RGBColor(0x44, 0x44, 0x44)
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x3E, 0x8C, 0x35)
NEAR_BLACK= RGBColor(0x1A, 0x1A, 0x1A)
CODE_BG   = RGBColor(0x1E, 0x1E, 0x1E)

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, font_size, bold,
             color_rgb, align=PP_ALIGN.LEFT, wrap=True, italic=False,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color_rgb
    run.font.name = font_name
    return tf


def add_top_bar(slide, title_text):
    """Red full-width top bar with white title text."""
    add_rect(slide, 0, 0, 13.33, 0.70, RED)
    add_text(slide, title_text, 0.25, 0.08, 12.8, 0.58,
             font_size=22, bold=True, color_rgb=WHITE, align=PP_ALIGN.LEFT)


def add_bullet_body(slide, bullets, top=0.85, left=0.35,
                    width=12.6, font_size=17, prefix="▶ ", color=DARK_GRAY):
    y = top
    for b in bullets:
        add_text(slide, prefix + b, left, y, width, 0.42,
                 font_size=font_size, bold=False, color_rgb=color,
                 align=PP_ALIGN.LEFT)
        y += 0.47


# ── Presentation setup ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)

BLANK = prs.slide_layouts[6]   # truly blank


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# Full dark background
add_rect(s1, 0, 0, 13.33, 7.50, NEAR_BLACK)

# Red left bar
add_rect(s1, 0, 0, 0.30, 7.50, RED)

# Title
add_text(s1, "Autonomous JBoss Operations\nwith Red Hat OpenShift AI",
         0.55, 2.20, 11.5, 1.80,
         font_size=40, bold=True, color_rgb=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text(s1,
         "Reducing MTTR and operational cost\nwith an AI-powered monitoring agent",
         0.55, 4.20, 10.0, 1.0,
         font_size=20, bold=False, color_rgb=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# Bottom-right confidential
add_text(s1, "Red Hat — Confidential",
         9.5, 7.05, 3.6, 0.35,
         font_size=12, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s2, "The Status Quo: Manual, Reactive, Expensive")

bullets = [
    "JBoss incidents go undetected until users complain",
    "Senior engineers spend 2+ hours per incident on root cause",
    "JIRA tickets are created manually — late, incomplete, inconsistent",
    "On-call rotations are expensive and cause burnout",
    "The same incidents recur — fixes never get documented",
]
add_bullet_body(s2, bullets, top=0.90, font_size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution (4 pillars)
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s3, "The JBoss AI Monitor: Detect. Analyse. Document. Automatically.")

box_w = 2.95
gap   = 0.13
start = 0.22
top_b = 1.45
h_b   = 3.60

pillars = [
    ("🔍", "Detect",      "Monitors pod state, logs, alerts and health every 60 seconds"),
    ("🧠", "Analyse",     "RHOAI-hosted LLM generates root cause + resolution steps"),
    ("🎫", "Document",    "Structured JIRA ticket created automatically with full context"),
    ("🔕", "Deduplicate", "Suppresses repeat tickets within a configurable time window"),
]

for i, (icon, title, body) in enumerate(pillars):
    x = start + i * (box_w + gap)

    # light gray bg
    add_rect(s3, x, top_b, box_w, h_b, LIGHT_GRAY)

    # red top border strip
    add_rect(s3, x, top_b, box_w, 0.07, RED)

    # icon
    add_text(s3, icon, x + 0.10, top_b + 0.18, box_w - 0.20, 0.75,
             font_size=38, bold=False, color_rgb=DARK_GRAY, align=PP_ALIGN.LEFT)

    # box title
    add_text(s3, title, x + 0.10, top_b + 1.00, box_w - 0.20, 0.40,
             font_size=16, bold=True, color_rgb=RED, align=PP_ALIGN.LEFT)

    # box body
    add_text(s3, body, x + 0.10, top_b + 1.50, box_w - 0.20, 1.90,
             font_size=13, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Time Savings
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s4, "2.5 Hours of Senior Engineer Time → 15 Minutes of Review")

col_x  = [0.25, 3.55, 7.20, 10.40]
col_w  = [3.20, 3.55, 3.10, 2.80]
rows = [
    (RED,                  WHITE,      ["Step",             "Before",          "After",        "Time Saved"]),
    (LIGHT_GRAY,           DARK_GRAY,  ["Issue detection",  "5–30 min",        "60 seconds",   "~29 min saved"]),
    (RGBColor(0xEB,0xEB,0xEB), DARK_GRAY, ["Root cause analysis","60–120 min",  "AI: 14 sec",   "~119 min saved"]),
    (LIGHT_GRAY,           DARK_GRAY,  ["JIRA ticket creation","10–15 min",    "Automated",    "~13 min saved"]),
    (NEAR_BLACK,           WHITE,      ["TOTAL",            "~2.5 hours",      "~15 min review","~2 hours saved"]),
]

row_h  = 0.62
top_r  = 0.80

for r, (bg, fg, cells) in enumerate(rows):
    y = top_r + r * row_h
    bold = (r == 0 or r == 4)
    for c, (cx, cw, cell) in enumerate(zip(col_x, col_w, cells)):
        add_rect(s4, cx, y, cw, row_h, bg)
        add_text(s4, cell, cx + 0.08, y + 0.10, cw - 0.15, row_h - 0.10,
                 font_size=14, bold=bold, color_rgb=fg, align=PP_ALIGN.LEFT)

# callout
callout_y = top_r + len(rows) * row_h + 0.15
add_rect(s4, 0.25, callout_y, 12.83, 0.55, RGBColor(0xFF, 0xEE, 0xEE))
add_text(s4, "At $100/hr blended SRE rate → $225 saved per incident",
         0.40, callout_y + 0.08, 12.50, 0.40,
         font_size=15, bold=True, color_rgb=RED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROI by Size
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s5, "The Value Scales With Your Environment")

tiers = [
    (LIGHT_GRAY,            DARK_GRAY, "Small",  "1 env · 10 incidents/mo",  "$27,000/yr saved"),
    (RGBColor(0xE8,0xE8,0xE8), DARK_GRAY, "Medium", "3 envs · 30 incidents/mo", "$81,000/yr saved"),
    (RED,                   WHITE,     "Large",  "10 envs · 100 incidents/mo","$270,000/yr saved"),
]
box_w5 = 4.0
gap5   = 0.22
start5 = 0.25
top5   = 1.10
h5     = 4.50

for i, (bg, fg, label, desc, amount) in enumerate(tiers):
    x = start5 + i * (box_w5 + gap5)
    add_rect(s5, x, top5, box_w5, h5, bg)

    # tier label
    add_text(s5, label, x+0.15, top5+0.25, box_w5-0.30, 0.55,
             font_size=22, bold=True, color_rgb=fg, align=PP_ALIGN.CENTER)

    # env desc
    add_text(s5, desc, x+0.15, top5+0.90, box_w5-0.30, 0.50,
             font_size=14, bold=False, color_rgb=fg, align=PP_ALIGN.CENTER)

    # dollar amount
    add_text(s5, amount, x+0.15, top5+1.60, box_w5-0.30, 0.80,
             font_size=32, bold=True, color_rgb=fg, align=PP_ALIGN.CENTER)

# footer
add_text(s5,
         "Conservative estimate — direct engineer time only. "
         "Does not include on-call premiums, prevented downtime, or revenue protection.",
         0.25, 6.30, 12.83, 0.60,
         font_size=12, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.CENTER,
         italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — On-Call Value
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s6, "24/7 Coverage — Without 3am Pages")

# quote box
add_rect(s6, 0.60, 0.90, 12.13, 2.40, NEAR_BLACK)
add_text(s6,
         '"The first P1 incident it catches at 3am on a Sunday\n'
         '— without waking anyone —\npays for the entire implementation."',
         0.80, 1.00, 11.73, 2.20,
         font_size=22, bold=False, italic=True,
         color_rgb=WHITE, align=PP_ALIGN.CENTER)

# stat boxes
stats = [
    ("10–25%", "On-call salary premium eliminated"),
    ("3am",    "Agent works nights, weekends, holidays"),
    ("0 min",  "Time to ticket after detection"),
]
sw = 3.90
sg = 0.22
sx = 0.60
sy = 3.55
sh = 2.70

for i, (num, label) in enumerate(stats):
    xx = sx + i*(sw+sg)
    add_rect(s6, xx, sy, sw, sh, LIGHT_GRAY)
    add_text(s6, num, xx+0.15, sy+0.55, sw-0.30, 0.90,
             font_size=36, bold=True, color_rgb=RED, align=PP_ALIGN.CENTER)
    add_text(s6, label, xx+0.15, sy+1.55, sw-0.30, 0.90,
             font_size=14, bold=False, color_rgb=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Sample JIRA Ticket
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s7, "Every Incident Gets a Structured, Actionable Ticket")

# dark code box
add_rect(s7, 0.25, 0.85, 9.30, 6.30, CODE_BG)

ticket_text = (
    "\U0001F534 OOMKilled — jboss-instance-0  |  CRITICAL  |  jboss-production\n\n"
    "\U0001F9E0 Root Cause (AI — Confidence: HIGH)\n"
    "JVM heap exhausted. -Xmx setting too low for workload.\n"
    "Container memory limit: 1536Mi. Current heap: 256Mi.\n\n"
    "\U0001F527 Resolution Steps\n"
    "1. Increase -Xmx to 1024m in standalone.xml\n"
    "2. Set OpenShift memory limit to 1536Mi\n"
    "3. Enable -XX:+HeapDumpOnOutOfMemoryError\n\n"
    "\U0001F6E1️ Prevention Tips\n"
    "• Set JVM heap to 75% of container memory limit\n"
    "• Add readiness probe on /health/ready\n\n"
    "\U0001F4DA References\n"
    "• Red Hat KBase: Tuning JVM for EAP on OpenShift"
)

add_text(s7, ticket_text,
         0.40, 0.95, 9.00, 6.10,
         font_size=10, bold=False, color_rgb=WHITE,
         align=PP_ALIGN.LEFT, font_name="Courier New")

# right callout
add_rect(s7, 9.80, 2.20, 3.28, 1.50, RGBColor(0xE8, 0xF5, 0xE6))
add_text(s7, "Created in 14 seconds\nNo human involved",
         9.90, 2.40, 3.08, 1.10,
         font_size=14, bold=True, color_rgb=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Red Hat Platform
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s8, "Runs on Technology You Already Own")

# Left table header
platform_rows = [
    (RED,       WHITE,      "Component",              "Red Hat Product"),
    (LIGHT_GRAY,DARK_GRAY,  "JBoss runtime",          "JBoss EAP / WildFly Operator"),
    (RGBColor(0xEB,0xEB,0xEB),DARK_GRAY,"Container platform","Red Hat OpenShift"),
    (LIGHT_GRAY,DARK_GRAY,  "AI inference",           "Red Hat OpenShift AI (RHOAI)"),
    (RGBColor(0xEB,0xEB,0xEB),DARK_GRAY,"Model serving","KServe + vLLM on GPU"),
    (LIGHT_GRAY,DARK_GRAY,  "Container image",        "UBI9 (Universal Base Image)"),
]

col1_x, col2_x = 0.25, 4.10
col1_w, col2_w = 3.75, 3.75
row_h8 = 0.70
top8   = 0.82

for r, (bg, fg, c1, c2) in enumerate(platform_rows):
    y = top8 + r * row_h8
    bold = (r == 0)
    add_rect(s8, col1_x, y, col1_w, row_h8, bg)
    add_rect(s8, col2_x, y, col2_w, row_h8, bg)
    add_text(s8, c1, col1_x+0.08, y+0.12, col1_w-0.15, row_h8-0.15,
             font_size=14, bold=bold, color_rgb=fg)
    add_text(s8, c2, col2_x+0.08, y+0.12, col2_w-0.15, row_h8-0.15,
             font_size=14, bold=bold, color_rgb=fg)

# Right callout boxes
callouts = [
    "✓  Activates your RHOAI investment",
    "✓  No new vendor relationship",
    "✓  Data never leaves your cluster",
]
right_x = 8.20
right_w = 4.88
cy = 1.10
for txt in callouts:
    add_rect(s8, right_x, cy, right_w, 0.85, LIGHT_GRAY)
    add_text(s8, txt, right_x+0.15, cy+0.12, right_w-0.30, 0.62,
             font_size=15, bold=True, color_rgb=DARK_GRAY)
    cy += 1.05


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Competitive Differentiation
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s9, "No Generic APM Tool Does This")

comp_cols_x = [0.25, 4.30, 7.40, 10.10]
comp_cols_w = [3.95, 3.00, 2.60, 3.00]
comp_rows = [
    ["Capability",                "Generic APM", "Cloud AI", "This Solution"],
    ["JBoss-specific root cause", "✗",           "Partial",  "✓"],
    ["AI resolution steps",       "✗",           "✗",        "✓"],
    ["On-premises LLM",           "✗",           "✗",        "✓"],
    ["Data stays in cluster",     "✗",           "✗",        "✓"],
    ["Structured JIRA tickets",   "✗",           "✗",        "✓"],
    ["No per-API-call cost",      "—",           "✗",        "✓"],
]

crow_h = 0.68
ctop   = 0.82

for r, row in enumerate(comp_rows):
    y = ctop + r * crow_h
    for c, (cx, cw, cell) in enumerate(zip(comp_cols_x, comp_cols_w, row)):
        if r == 0:
            bg = RED if c == 3 else DARK_GRAY
            fg = WHITE
            bold = True
        elif r % 2 == 0:
            bg = RGBColor(0xEB, 0xEB, 0xEB)
            fg = DARK_GRAY
            bold = False
        else:
            bg = LIGHT_GRAY
            fg = DARK_GRAY
            bold = False

        # colour ✓ and ✗ differently
        if cell == "✓" and r > 0:
            fg = GREEN
            bold = True
        elif cell == "✗" and r > 0:
            fg = MID_GRAY
            bold = False

        add_rect(s9, cx, y, cw, crow_h, bg)
        add_text(s9, cell, cx+0.08, y+0.10, cw-0.15, crow_h-0.12,
                 font_size=14, bold=bold, color_rgb=fg,
                 align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Time to Value
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s10, "Running in Your Environment in Under a Day")

phases = [
    ("Phase 1", "2–4 hrs", "Deploy",
     "OpenShift manifests, RHOAI endpoint, JIRA config"),
    ("Phase 2", "1–2 hrs", "Validate",
     "Inject test errors, confirm ticket creation, tune patterns"),
    ("Phase 3", "Ongoing", "Expand",
     "Add namespaces, extend to Quarkus/Spring Boot"),
]

pw = 3.80
pg = 0.27
px = 0.28
py = 1.20
ph = 4.50

for i, (phase, duration, title, body) in enumerate(phases):
    xx = px + i*(pw+pg)
    add_rect(s10, xx, py, pw, ph, LIGHT_GRAY)
    add_rect(s10, xx, py, pw, 0.07, RED)

    add_text(s10, phase, xx+0.15, py+0.18, pw-0.30, 0.40,
             font_size=13, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s10, duration, xx+0.15, py+0.62, pw-0.30, 0.45,
             font_size=20, bold=True, color_rgb=RED, align=PP_ALIGN.CENTER)
    add_text(s10, title, xx+0.15, py+1.18, pw-0.30, 0.50,
             font_size=18, bold=True, color_rgb=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_text(s10, body, xx+0.15, py+1.80, pw-0.30, 2.40,
             font_size=13, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.CENTER)

    # arrow connector (except after last)
    if i < 2:
        ax = xx + pw + 0.04
        add_text(s10, "→", ax, py + 1.80, pg + 0.02, 0.60,
                 font_size=22, bold=True, color_rgb=RED, align=PP_ALIGN.CENTER)

# bottom note
add_text(s10,
         "No professional services engagement required for initial deployment.",
         0.25, 6.00, 12.83, 0.55,
         font_size=13, bold=False, italic=True,
         color_rgb=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Call to Action
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_rect(s11, 0, 0, 13.33, 7.50, WHITE)
add_top_bar(s11, "Three Ways to Move Forward Today")

options = [
    ("Option 1 — Proof of Value", "1 week",
     "Deploy in non-production. Run for one week. Measure MTTR delta vs baseline."),
    ("Option 2 — Architecture Workshop", "Half day",
     "Red Hat architect maps solution to your JBoss topology and RHOAI setup."),
    ("Option 3 — Production Deployment", "Scoped engagement",
     "Deploy, configure, hand over with runbook and training."),
]

ow = 3.85
og = 0.22
ox = 0.28
oy = 0.90
oh = 4.90

for i, (title, duration, body) in enumerate(options):
    xx = ox + i*(ow+og)
    add_rect(s11, xx, oy, ow, oh, LIGHT_GRAY)
    add_rect(s11, xx, oy, ow, 0.07, RED)

    add_text(s11, title, xx+0.15, oy+0.20, ow-0.30, 0.60,
             font_size=14, bold=True, color_rgb=DARK_GRAY, align=PP_ALIGN.LEFT)
    add_text(s11, duration, xx+0.15, oy+0.88, ow-0.30, 0.42,
             font_size=14, bold=True, color_rgb=RED, align=PP_ALIGN.LEFT)
    add_text(s11, body, xx+0.15, oy+1.40, ow-0.30, 3.20,
             font_size=13, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.LEFT)

# footer
add_text(s11,
         "All source code is open and customer-owned.",
         0.25, 6.20, 12.83, 0.55,
         font_size=13, bold=False, italic=True,
         color_rgb=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You / Contact
# ══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)

add_rect(s12, 0, 0, 13.33, 7.50, NEAR_BLACK)
add_rect(s12, 0, 0, 0.30, 7.50, RED)

add_text(s12, "Questions?",
         0.55, 2.30, 11.5, 1.20,
         font_size=52, bold=True, color_rgb=WHITE, align=PP_ALIGN.LEFT)

add_text(s12, "Let's talk about your JBoss environment.",
         0.55, 3.75, 10.0, 0.70,
         font_size=22, bold=False, color_rgb=LIGHT_GRAY, align=PP_ALIGN.LEFT)

add_text(s12,
         "[Name]  ·  [email]  ·  Red Hat Senior Architect",
         0.55, 6.70, 11.5, 0.55,
         font_size=14, bold=False, color_rgb=MID_GRAY, align=PP_ALIGN.LEFT)


# ── Save ───────────────────────────────────────────────────────────────────────
out_path = "/home/user/jboss-ai-agent-monitor/docs/jboss-ai-monitor-deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
